"""
Build a schematic slide deck synthesizing the four meetings in
TRANSCRIPT FROM MEETINGS.pdf into the new direction of the ED fire
optimization project.

Output: meeting_synthesis.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

NAVY = RGBColor(0x0B, 0x2A, 0x4A)
RED  = RGBColor(0xC0, 0x39, 0x2B)
ORG  = RGBColor(0xE6, 0x7E, 0x22)
GRN  = RGBColor(0x1E, 0x8A, 0x4A)
GRAY = RGBColor(0x55, 0x55, 0x55)
LITE = RGBColor(0xEC, 0xEF, 0xF1)
WHT  = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_title(slide, text, sub=None):
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.9))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = NAVY
    if sub:
        p2 = tf.add_paragraph(); p2.text = sub
        p2.font.size = Pt(14); p2.font.color.rgb = GRAY


def box(slide, x, y, w, h, text, fill=NAVY, fg=WHT, size=14, bold=True, align=PP_ALIGN.CENTER):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = fill
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = fg
    return s


def text(slide, x, y, w, h, lines, size=12, color=NAVY, bold=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    if isinstance(lines, str): lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln; p.font.size = Pt(size); p.font.color.rgb = color; p.font.bold = bold
    return tb


def arrow(slide, x1, y1, x2, y2, color=GRAY):
    c = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(2.25)
    try:
        from pptx.oxml.ns import qn
        ln = c.line._get_or_add_ln()
        from lxml import etree
        tail = etree.SubElement(ln, qn('a:tailEnd'))
        tail.set('type', 'triangle')
    except Exception:
        pass
    return c


# ---------- Slide 1: Title ----------
s = prs.slides.add_slide(BLANK)
add_title(s, "ED Fire Optimization — Meeting Synthesis & New Direction",
          "Mapping 4 meetings into the revised research arc · prepared for sync between Mac & Windows workstations")
box(s, 1.0, 1.8, 11.3, 1.0, "Old: hand-built Model C, rank #1 offline (0.6713).  New: dual-track optimizer + multi-candidate fire model + 3-chapter dissertation under a fire+carbon-climate trajectory.",
    fill=LITE, fg=NAVY, size=16, bold=False)
box(s, 1.0, 3.1, 5.4, 0.7, "Meeting 1 — Strategic / Dissertation framing", fill=NAVY)
box(s, 6.9, 3.1, 5.4, 0.7, "Meeting 2 — Model C reproduction milestones", fill=ORG)
box(s, 1.0, 4.0, 5.4, 0.7, "Meeting 3 — Dev joins, optimization architecture", fill=GRN)
box(s, 6.9, 4.0, 5.4, 0.7, "Meeting 4 — Literature + process-vs-empirical", fill=RED)
text(s, 1.0, 5.2, 11.3, 1.5, [
    "Read in order. Slides 3–6 expand each meeting; slide 7 fuses them into the new direction; slide 8 lists open items.",
    "Companion file: PROGRESS.md (repo root) — where the actual code/data state lives day-to-day.",
], size=13, color=GRAY)

# ---------- Slide 2: 30,000-ft view of the dissertation arc ----------
s = prs.slides.add_slide(BLANK)
add_title(s, "30,000-ft view: the new dissertation arc",
          "Trajectory: fire × carbon-climate (past · present · near-future · long-term)")
box(s, 0.5, 1.6, 3.9, 1.4,
    "CHAPTER 1\nImproved offline ED fire model\n(THIS project — Model C → multi-candidate)",
    fill=NAVY, size=15)
box(s, 4.7, 1.6, 3.9, 1.4,
    "CHAPTER 2\nNear-future seasonal forecasts\n(El Niño / La Niña → regional fire risk, weeks–years out)",
    fill=ORG, size=15)
box(s, 8.9, 1.6, 3.9, 1.4,
    "CHAPTER 3\nLong-term scenarios (SSPs)\n(coupled ESM runs, 100-yr fire futures)",
    fill=RED, size=15)
arrow(s, 4.4, 2.3, 4.7, 2.3); arrow(s, 8.6, 2.3, 8.9, 2.3)

box(s, 0.5, 3.4, 12.3, 0.7,
    "Common throughline:  improved fire model  →  near-term value (budget, mitigation, prescribed burns)  →  long-term value (carbon-climate projections)",
    fill=LITE, fg=NAVY, size=13, bold=False)

text(s, 0.5, 4.3, 12.3, 0.4,
    "WHY US (Venn intersection from Meeting 1):", size=14, bold=True)
box(s, 0.5, 4.7, 2.9, 0.7, "Skills / interest", fill=GRN, size=12)
box(s, 3.6, 4.7, 2.9, 0.7, "Application (NASA FINESST)", fill=GRN, size=12)
box(s, 6.7, 4.7, 2.9, 0.7, "Timeliness (ENSO + climate forcing)", fill=GRN, size=12)
box(s, 9.8, 4.7, 3.0, 0.7, "Uniquely qualified", fill=GRN, size=12)
text(s, 0.5, 5.6, 12.3, 1.4, [
    "Goal phrase: NEW · ADVANCED · IMPORTANT · USEFUL.",
    "Deliverable by end of semester:  pre-proposal (5 pages) + annotated bibliography → NASA FINESST format proposal next year.",
    "Defend-the-question test:  \"Why does the world need another fire-and-ecosystems study?\" — must have a sharp 1-sentence answer.",
], size=13, color=GRAY)

# ---------- Slide 3: Meeting 1 — strategic ----------
s = prs.slides.add_slide(BLANK)
add_title(s, "Meeting 1 — Strategic framing (with George)",
          "Setting up the dissertation around fire; sharpening the science question")
text(s, 0.5, 1.3, 6.0, 0.4, "Sharpening the science question:", size=14, bold=True)
box(s, 0.5, 1.7, 6.0, 0.6, "v1: \"What is the role of fire in carbon-climate?\"  (too big — career-sized)",
    fill=LITE, fg=NAVY, bold=False, size=12, align=PP_ALIGN.LEFT)
box(s, 0.5, 2.4, 6.0, 0.6, "v2: \"…over the United States, 2020–2025?\"  (more tractable — but: carbon? albedo? latent heat?)",
    fill=LITE, fg=NAVY, bold=False, size=12, align=PP_ALIGN.LEFT)
box(s, 0.5, 3.1, 6.0, 0.7, "v3 (target):  Carbon-cycle effect of US fires, 2020–2025 — biogeochemical & biophysical separated",
    fill=NAVY, size=13, align=PP_ALIGN.LEFT)

text(s, 0.5, 4.0, 6.0, 0.4, "Three candidate trajectories considered:", size=14, bold=True)
text(s, 0.5, 4.4, 6.0, 2.6, [
    "• Carbon-only  (past · present · near-future · long-term)",
    "• Carbon + biodiversity",
    "• Carbon + ecosystem-services / human health",
    "→ Chosen direction: stay in carbon, but add the predictive dimension",
    "   (chapter 1 = credible model, chapter 2 = near-future, chapter 3 = scenarios)",
], size=12, color=GRAY)

text(s, 7.0, 1.3, 5.8, 0.4, "Why-Us Venn (the test George keeps coming back to):", size=14, bold=True)
box(s, 7.5, 1.9, 4.8, 0.6, "Skills / interest (GIS, comms, ED model exposure)", fill=GRN, size=11)
box(s, 7.5, 2.6, 4.8, 0.6, "Application location (US — managed, data-rich)", fill=GRN, size=11)
box(s, 7.5, 3.3, 4.8, 0.6, "Timeliness (NASA reorganizing, FINESST cycle)", fill=GRN, size=11)
box(s, 7.5, 4.0, 4.8, 0.6, "Technical advantage (ED + remote sensing link)", fill=GRN, size=11)
box(s, 7.5, 4.9, 4.8, 0.7, "→ Intersection = the dissertation's defensible niche",
    fill=NAVY, size=13)
text(s, 7.0, 5.9, 5.8, 1.0,
    "Operational rule: anything that fails any one circle → drop it.",
    size=12, color=GRAY, bold=True)

# ---------- Slide 4: Meeting 2 — Model C reproduction ----------
s = prs.slides.add_slide(BLANK)
add_title(s, "Meeting 2 — Model C reproduction milestones",
          "Three milestones; we are between M2 and M3")
box(s, 0.5, 1.5, 4.0, 1.4,
    "M1 — Understand\nWhat ED actually reports for fire / burned area in TRENDY",
    fill=GRN, size=14)
box(s, 4.7, 1.5, 4.0, 1.4,
    "M2 — Reproduce\nOffline pipeline matches ED's reported numbers (close enough)",
    fill=ORG, size=14)
box(s, 8.9, 1.5, 4.0, 1.4,
    "M3 — Improve\nFit to GFED, optimize parameters, evaluate on hold-out",
    fill=RED, size=14)
arrow(s, 4.5, 2.2, 4.7, 2.2); arrow(s, 8.7, 2.2, 8.9, 2.2)

text(s, 0.5, 3.2, 12.3, 0.4, "Inputs assembled offline (CRUJRA v3.5 + ED outputs):", size=14, bold=True)
box(s, 0.5, 3.7, 2.3, 0.7, "PET (temp-only, simple)", fill=NAVY, size=11)
box(s, 3.0, 3.7, 2.3, 0.7, "Dryness index D̄", fill=NAVY, size=11)
box(s, 5.5, 3.7, 2.3, 0.7, "Ignition factor", fill=NAVY, size=11)
box(s, 8.0, 3.7, 2.3, 0.7, "Aboveground biomass (from ED)", fill=NAVY, size=11)
box(s, 10.5, 3.7, 2.3, 0.7, "Burned area (target = GFED)", fill=NAVY, size=11)

text(s, 0.5, 4.8, 12.3, 0.4, "Open issues flagged in this meeting:", size=14, bold=True)
text(s, 0.5, 5.2, 12.3, 2.0, [
    "• PET formula:  used a temperature-only proxy because full Penman-Monteith inputs not all in offline drivers — must verify.",
    "• Precip units:  data labelled mm/6 h but actually monthly — confirmed; otherwise dryness index goes to zero.",
    "• Biomass-fire response is too steep / nearly binary — burns whole grid cells or nothing.  Want a smoother distribution.",
    "• Calibration window:  2001–2016 (ILAMB evaluation period).  Hold out 2024 as an out-of-sample test.",
    "• Mean-of-fire ≠ fire-of-mean:  must compute fire each year and then average, not the other way around.",
    "• GFED4 vs GFED5:  GFED5 has more burned area (more small fires) but similar carbon emissions — Giglio & Holiday are in our dept.",
], size=12, color=GRAY)

# ---------- Slide 5: Meeting 3 — optimization architecture ----------
s = prs.slides.add_slide(BLANK)
add_title(s, "Meeting 3 — Dev joins · the new optimization architecture",
          "Outer optimizer wraps the offline fire model · two parallel methods · multi-candidate output")

# Big architecture diagram
box(s, 0.5, 1.4, 2.4, 1.0, "Optimizer\n(outer loop)", fill=NAVY, size=13)
box(s, 3.6, 1.4, 2.6, 1.0, "Offline\nfire model\nf(params, drivers)", fill=ORG, size=13)
box(s, 6.9, 1.4, 2.4, 1.0, "Goodness-of-fit\n(SSE / likelihood)", fill=GRN, size=13)
box(s, 10.0, 1.4, 2.8, 1.0, "Benchmark\nGFED4 / GFED5", fill=RED, size=13)
arrow(s, 2.9, 1.9, 3.6, 1.9)            # opt -> model
arrow(s, 6.2, 1.9, 6.9, 1.9)            # model -> gof
arrow(s, 10.0, 1.9, 9.3, 1.9, color=RED)# bench -> gof
# feedback
arrow(s, 8.1, 2.4, 8.1, 3.0, color=NAVY)
arrow(s, 8.1, 3.0, 1.7, 3.0, color=NAVY)
arrow(s, 1.7, 3.0, 1.7, 2.4, color=NAVY)
text(s, 1.8, 3.05, 6.0, 0.3, "params update ←", size=10, color=NAVY, bold=True)

# Two tracks
text(s, 0.5, 3.6, 12.3, 0.4, "Two parallel optimization tracks (run BOTH, compare):", size=14, bold=True)
box(s, 0.5, 4.1, 6.1, 1.2,
    "Track A — Dev (CS-style)\nFully automated; tries many functional forms;\nempirical-leaning.  Fast, broad search.",
    fill=GRN, size=12)
box(s, 6.7, 4.1, 6.1, 1.2,
    "Track B — Manual / physical\nFix functional form on physical grounds;\noptimize parameter values within it.",
    fill=NAVY, size=12)

# Outputs
box(s, 0.5, 5.6, 12.3, 0.7,
    "Output:  multiple CANDIDATE fire models (best, simpler-but-close, complex-but-explainable, etc.) → team picks which to drop into ED",
    fill=LITE, fg=NAVY, bold=False, size=13)

text(s, 0.5, 6.5, 12.3, 0.8, [
    "New mechanism on the table:  add canopy/biomass HEIGHT as a 3rd predictor (alongside dryness & biomass) — only below-canopy biomass is flammable.",
    "Debug ask from George:  print D̄ at every step in the offline code so we know exactly what value the formula is consuming.",
], size=12, color=GRAY)

# ---------- Slide 6: Meeting 4 — process vs empirical ----------
s = prs.slides.add_slide(BLANK)
add_title(s, "Meeting 4 — Literature + process-vs-empirical philosophy",
          "Why the multi-candidate output matters")

# Continuum bar
box(s, 0.5, 1.5, 2.6, 0.8, "Pure empirical\n(black-box AI)", fill=RED, size=12)
box(s, 3.3, 1.5, 6.7, 0.8, "← continuum →   most real models live here", fill=LITE, fg=NAVY, bold=False, size=12)
box(s, 10.2, 1.5, 2.6, 0.8, "Pure process-based\n(every term derived)", fill=GRN, size=12)
text(s, 0.5, 2.4, 12.3, 0.4, "Guiding principle:  \"best prediction for the right reason.\"", size=14, bold=True, color=NAVY)

# Two columns
text(s, 0.5, 3.0, 6.0, 0.4, "Papers / models discussed:", size=14, bold=True)
text(s, 0.5, 3.4, 6.0, 3.5, [
    "• Cardoso (Amazon, multiplicative form):",
    "    fire = fuel × availability × flammability × ignition × …",
    "    — if any factor = 0, fire = 0.",
    "    — Pete-Bruce road-building added human ignition access.",
    "• Manuel Cardoso:  lightning + human-activity ignition layer.",
    "• Manwald (2002, US):  added a fire-suppression factor that ramps with USDA spending.",
    "• Most existing DGVMs treat fire as fixed parameter, not interactive.",
], size=12, color=GRAY)

text(s, 7.0, 3.0, 5.8, 0.4, "Goodness-of-fit:  move toward likelihood-based.", size=14, bold=True)
text(s, 7.0, 3.4, 5.8, 3.5, [
    "• χ² ≈ Σ (model − data)² / σ²    (then ÷ 2σ², minus log term)",
    "• Likelihood = how likely is THIS dataset under the model.",
    "• Optimizer choice:",
    "    – Gradient-based:  fast, gets stuck in local optima.",
    "    – Metropolis (MCMC): slower, escapes local optima,",
    "       gives uncertainty as a bonus.",
    "• Plan: try both; report which candidate each one finds.",
], size=12, color=GRAY)

# ---------- Slide 7: New direction — synthesis ----------
s = prs.slides.add_slide(BLANK)
add_title(s, "Synthesis — what changed about this project",
          "Old single-model, leaderboard-chasing run → multi-candidate, statistically-grounded chapter 1")

# Old vs New table
box(s, 0.5, 1.4, 6.1, 0.7, "BEFORE", fill=GRAY, size=14)
box(s, 6.7, 1.4, 6.1, 0.7, "AFTER (new direction)", fill=NAVY, size=14)

rows = [
    ("One hand-tuned formula (Model C)", "Multiple candidate models, picked by team"),
    ("Optuna over 12 params, single run", "Two parallel optimizers (Dev's CS + manual physical)"),
    ("Inputs: dryness, precip, GPP, T_air", "Above + height (canopy / below-canopy biomass split)"),
    ("Score = ILAMB Overall (offline)", "Score = likelihood vs GFED + held-out 2024"),
    ("Endpoint = leaderboard #1", "Endpoint = chapter 1 of dissertation + NASA FINESST"),
    ("Standalone deliverable", "Feeds chapter 2 (seasonal) & chapter 3 (scenarios)"),
]
for i, (a, b) in enumerate(rows):
    y = 2.2 + i * 0.65
    box(s, 0.5, y, 6.1, 0.55, a, fill=LITE, fg=NAVY, bold=False, size=12, align=PP_ALIGN.LEFT)
    box(s, 6.7, y, 6.1, 0.55, b, fill=LITE, fg=NAVY, bold=False, size=12, align=PP_ALIGN.LEFT)

text(s, 0.5, 6.3, 12.3, 1.0, [
    "One-line summary:  Model C becomes a baseline, not the goal.  We're building the optimization machinery that produces a family of fire models,",
    "then defending the chosen one on physical + statistical grounds — and that becomes Chapter 1 of the dissertation.",
], size=13, color=NAVY, bold=True)

# ---------- Slide 8: Open items / next steps ----------
s = prs.slides.add_slide(BLANK)
add_title(s, "Open items & immediate next steps",
          "What to do this week / before the next meeting")

text(s, 0.5, 1.3, 6.0, 0.4, "Technical (code & data):", size=14, bold=True)
text(s, 0.5, 1.7, 6.0, 4.5, [
    "1. Add print/log of D̄ everywhere it's computed in the offline pipeline,",
    "    so George can sanity-check the values.",
    "2. Wire up an outer optimizer interface:  fire_model(params, drivers) → SSE,",
    "    callable as a function.  Coordinate API with Dev.",
    "3. Re-segment train/test:  calibrate on 2001–2016, hold out 2024.",
    "4. Replace SSE with χ²/likelihood goodness-of-fit (carry σ²).",
    "5. Add HEIGHT as a candidate predictor; restrict flammable biomass",
    "    to below a threshold height.",
    "6. Swap GFED4 → GFED5 as a sensitivity check (talk to Giglio / Holiday).",
    "7. Maps of every input layer (PET, D̄, biomass, GPP, fire) — make sure they",
    "    pass the eyeball test (PET high in tropics, low in high-lat, etc.).",
], size=12, color=GRAY)

text(s, 7.0, 1.3, 5.8, 0.4, "Strategic (writing & framing):", size=14, bold=True)
text(s, 7.0, 1.7, 5.8, 4.5, [
    "A. Draft the 1-sentence answer to:  \"Why does the world need",
    "    another fire study?\"  (defend at next meeting.)",
    "B. Pre-proposal outline (~5 pages):  background · science questions ·",
    "    objectives · expected contribution.",
    "C. Annotated bibliography of papers covered (Cardoso, Manuel-Cardoso,",
    "    Manwald, etc.) — re-organized around chapter 1's argument.",
    "D. Track NASA reorganization → align format with FINESST.",
    "E. Schedule joint meeting:  Richard + George + Dev + Mike,",
    "    formalize optimization API and division of labor.",
], size=12, color=GRAY)

box(s, 0.5, 6.5, 12.3, 0.7,
    "Where the work is logged:  PROGRESS.md (repo root, syncs Mac ↔ Windows via Drive).  Update at end of every working session.",
    fill=NAVY, size=13)

prs.save("meeting_synthesis.pptx")
print("wrote meeting_synthesis.pptx")
