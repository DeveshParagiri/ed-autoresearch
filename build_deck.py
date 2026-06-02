"""
Single source of truth for tables/schematics for the ED fire optimization
project. Re-run any time something changes.

Output: figures_and_tables.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Palette — restrained, two-tone
INK   = RGBColor(0x1A, 0x1A, 0x1A)
SUBT  = RGBColor(0x6B, 0x72, 0x80)
RULE  = RGBColor(0xD0, 0xD5, 0xDB)
HEAD  = RGBColor(0x11, 0x2D, 0x4E)
ACCENT= RGBColor(0xC0, 0x39, 0x2B)
SOFT  = RGBColor(0xF4, 0xF6, 0xF8)
WHT   = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide_header(s, kicker, title):
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.text = kicker.upper()
    p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = ACCENT
    p.font.name = "Helvetica"

    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12.3), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = HEAD
    p.font.name = "Helvetica"

    line = s.shapes.add_connector(1, Inches(0.5), Inches(1.4), Inches(12.83), Inches(1.4))
    line.line.color.rgb = RULE; line.line.width = Pt(0.75)


def cell(tbl, r, c, text, *, bold=False, fg=INK, bg=None,
         align=PP_ALIGN.LEFT, size=11, anchor=MSO_ANCHOR.MIDDLE):
    cell = tbl.cell(r, c)
    cell.text = ""
    if bg is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
    cell.vertical_anchor = anchor
    cell.margin_left = cell.margin_right = Inches(0.08)
    cell.margin_top = cell.margin_bottom = Inches(0.04)
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    r0 = p.runs[0] if p.runs else p.add_run()
    if not p.runs: r0.text = text
    for run in p.runs:
        run.font.size = Pt(size); run.font.bold = bold
        run.font.color.rgb = fg; run.font.name = "Helvetica"


# ============================================================
# Slide 1 — Title
# ============================================================
s = prs.slides.add_slide(BLANK)
tb = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.7), Inches(0.5))
p = tb.text_frame.paragraphs[0]
p.text = "ED FIRE OPTIMIZATION"
p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = ACCENT
p.font.name = "Helvetica"

tb = s.shapes.add_textbox(Inches(0.8), Inches(2.9), Inches(11.7), Inches(1.4))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Advances in remote sensing alter the structural form of fire dynamics"
p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = HEAD
p.font.name = "Helvetica"

tb = s.shapes.add_textbox(Inches(0.8), Inches(4.4), Inches(11.7), Inches(0.6))
p = tb.text_frame.paragraphs[0]
p.text = "Working title (Hurtt, 2026-05-05). Tables & schematics live here. PROGRESS.md is the working log."
p.font.size = Pt(14); p.font.color.rgb = SUBT
p.font.name = "Helvetica"

# ============================================================
# Slide 2 — The 6-row evaluation table
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Table 1",
             "Six fire models · 3 optimization criteria × 2 reference datasets")

# Table dims: header row + 6 data rows; columns: #, Calib, Opt criterion, Selected inputs, Global, By continent, By fire type, Global, By continent, By fire type
n_rows = 1 + 6 + 1   # +1 for super-header grouping
n_cols = 10

tbl_shape = s.shapes.add_table(n_rows, n_cols,
                                Inches(0.5), Inches(1.7),
                                Inches(12.33), Inches(4.6))
tbl = tbl_shape.table

# Column widths (sum = 12.33)
widths = [0.4, 1.1, 2.0, 2.4, 0.85, 1.05, 1.05, 0.85, 1.05, 1.05]
for i, w in enumerate(widths):
    tbl.columns[i].width = Inches(w)

# Row heights
tbl.rows[0].height = Inches(0.30)   # super header
tbl.rows[1].height = Inches(0.55)   # main header
for r in range(2, n_rows):
    tbl.rows[r].height = Inches(0.55)

# Super-header row (groups)
for c in range(4):
    cell(tbl, 0, c, "", bg=WHT)
cell(tbl, 0, 4, "Evaluated on GFED4", bold=True, fg=WHT, bg=HEAD,
     align=PP_ALIGN.CENTER, size=11)
# merge 4..6
tbl.cell(0, 4).merge(tbl.cell(0, 6))
cell(tbl, 0, 7, "Evaluated on GFED5", bold=True, fg=WHT, bg=ACCENT,
     align=PP_ALIGN.CENTER, size=11)
tbl.cell(0, 7).merge(tbl.cell(0, 9))

# Main header row
headers = ["#", "Calib. dataset", "Optimization criterion",
           "Inputs (expected)",
           "Global R²", "By continent\n(R², n=5–6)", "By fire type\n(R², n=4–5)",
           "Global R²", "By continent\n(R², n=5–6)", "By fire type\n(R², n=4–5)"]
for c, h in enumerate(headers):
    cell(tbl, 1, c, h, bold=True, fg=HEAD, bg=SOFT,
         align=PP_ALIGN.CENTER, size=10)

# Data rows
rows = [
    ("1", "GFED4", "Opt 1 — every fire equal",        "Dryness, GPP   (= current Model C)"),
    ("2", "GFED4", "Opt 2 — every continent equal",   "Dryness, GPP  (+ Fuel?)"),
    ("3", "GFED4", "Opt 3 — every fire type equal",   "Dryness, GPP, Fuel"),
    ("4", "GFED5", "Opt 1 — every fire equal",        "+ Ignition / Land-use?"),
    ("5", "GFED5", "Opt 2 — every continent equal",   "+ Ignition / Land-use?"),
    ("6", "GFED5", "Opt 3 — every fire type equal",   "Dryness, GPP, Fuel, Ignition, LU mask"),
]
for i, row in enumerate(rows):
    r = 2 + i
    bg = WHT if i % 2 == 0 else SOFT
    cell(tbl, r, 0, row[0], bold=True, fg=HEAD, bg=bg, align=PP_ALIGN.CENTER, size=11)
    cell(tbl, r, 1, row[1], bg=bg, align=PP_ALIGN.CENTER, size=11)
    cell(tbl, r, 2, row[2], bg=bg, size=11)
    cell(tbl, r, 3, row[3], bg=bg, size=10, fg=SUBT)
    for c in range(4, 10):
        cell(tbl, r, c, "—", bg=bg, fg=SUBT, align=PP_ALIGN.CENTER, size=11)

# Footer caption
tb = s.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.9))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = ("Each row = one independently-optimized fire model. Each row gets evaluated against BOTH datasets "
          "(calibrated on one, cross-scored on the other) at three aggregations: global, by continent (~5–6), and by fire type (~4–5).")
p.font.size = Pt(11); p.font.color.rgb = SUBT; p.font.name = "Helvetica"
p2 = tf.add_paragraph()
p2.text = ("The paper's claim: cells in the “Inputs” column change across rows — i.e. structural change, not just parameter change.")
p2.font.size = Pt(11); p2.font.color.rgb = HEAD; p2.font.bold = True; p2.font.name = "Helvetica"

# ============================================================
# Slide 3 — Two axes schematic
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Schematic 1", "Two axes that produce the six models")

# Vertical axis label
tb = s.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(2.0), Inches(0.4))
p = tb.text_frame.paragraphs[0]; p.text = "Reference dataset ↓"
p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = HEAD; p.font.name = "Helvetica"

# Horizontal axis label
tb = s.shapes.add_textbox(Inches(2.7), Inches(1.9), Inches(8.0), Inches(0.4))
p = tb.text_frame.paragraphs[0]; p.text = "Optimization criterion →"
p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = HEAD; p.font.name = "Helvetica"

# 2x3 grid of cards
def card(s, x, y, w, h, title, sub, fill=SOFT, fg=INK):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = fill
    box.line.color.rgb = RULE; box.line.width = Pt(0.75)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.12)
    tf.margin_top = tf.margin_bottom = Inches(0.08)
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = fg; p.font.name = "Helvetica"
    p2 = tf.add_paragraph(); p2.text = sub
    p2.font.size = Pt(10); p2.font.color.rgb = SUBT; p2.font.name = "Helvetica"

# Row labels
tb = s.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(2.0), Inches(1.0))
p = tb.text_frame.paragraphs[0]; p.text = "GFED4"
p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = HEAD; p.font.name = "Helvetica"
p2 = tb.text_frame.add_paragraph(); p2.text = "current standard"
p2.font.size = Pt(10); p2.font.color.rgb = SUBT; p2.font.name = "Helvetica"

tb = s.shapes.add_textbox(Inches(0.5), Inches(4.7), Inches(2.0), Inches(1.0))
p = tb.text_frame.paragraphs[0]; p.text = "GFED5"
p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = ACCENT; p.font.name = "Helvetica"
p2 = tb.text_frame.add_paragraph(); p2.text = "~2× burned area"
p2.font.size = Pt(10); p2.font.color.rgb = SUBT; p2.font.name = "Helvetica"

# 6 cards
xs = [2.7, 6.1, 9.5]
titles = ["Opt 1 · every fire equal", "Opt 2 · every continent equal", "Opt 3 · every fire type equal"]
subs_4 = ["Dryness, GPP  (Model C)", "+ Fuel?", "Dryness, GPP, Fuel"]
subs_5 = ["+ Ignition / LU?", "+ Ignition / LU?", "Dryness, GPP, Fuel, Ignition, LU"]
for i, x in enumerate(xs):
    card(s, x, 2.6, 3.2, 1.5, titles[i], subs_4[i])
for i, x in enumerate(xs):
    card(s, x, 4.6, 3.2, 1.5, titles[i], subs_5[i],
         fill=RGBColor(0xFD, 0xEE, 0xEC))

# Note
tb = s.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = ("Hurtt's hypothesis: moving rightward (richer optimization criterion) and downward (richer dataset) "
          "forces the optimizer to keep additional input variables. That structural shift IS the paper's finding.")
p.font.size = Pt(11); p.font.color.rgb = SUBT; p.font.name = "Helvetica"

# ============================================================
# Slide 4 — Input ladder
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Schematic 2", "Input ladder · what the optimizer might select")

# Ladder of bigger boxes
levels = [
    ("Level 0", "Dryness", "moisture-driven combustibility"),
    ("Level 1", "+ GPP", "vegetation productivity (proxy for grass vs desert)"),
    ("Level 2", "+ Fuel", "fuel load — needed to lift overfit on grasslands"),
    ("Level 3", "+ Ignition", "humans / lightning — required by GFED5 small fires"),
    ("Level 4", "+ Land-use mask", "explicit cropland separation"),
]
y = 1.8
for i, (lvl, name, sub) in enumerate(levels):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(1.5), Inches(y), Inches(10.3), Inches(0.85))
    box.fill.solid(); box.fill.fore_color.rgb = SOFT if i % 2 == 0 else WHT
    box.line.color.rgb = RULE; box.line.width = Pt(0.75)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.18); tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = f"{lvl}    {name}    —  {sub}"
    p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = HEAD
    p.font.name = "Helvetica"
    y += 1.0

# Arrow on the left
arr = s.shapes.add_connector(1, Inches(1.0), Inches(1.9), Inches(1.0), Inches(6.6))
arr.line.color.rgb = ACCENT; arr.line.width = Pt(2.5)

tb = s.shapes.add_textbox(Inches(0.3), Inches(3.6), Inches(0.7), Inches(0.7))
p = tb.text_frame.paragraphs[0]
p.text = "richer →"
p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = ACCENT; p.font.name = "Helvetica"

# Footer
tb = s.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4))
p = tb.text_frame.paragraphs[0]
p.text = "Today's Model C sits at Level 1. The paper's hypothesis: GFED5 + Opt 3 lands the optimizer somewhere near Level 3–4."
p.font.size = Pt(11); p.font.color.rgb = SUBT; p.font.name = "Helvetica"

# ============================================================
# Slide 5 — Two deadlines
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Plan", "Two deliverables on different clocks")

def big_card(s, x, y, w, h, kicker, title, body, accent):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                              Inches(0.12), Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
    tb = s.shapes.add_textbox(Inches(x + 0.35), Inches(y + 0.1),
                                Inches(w - 0.5), Inches(0.4))
    p = tb.text_frame.paragraphs[0]; p.text = kicker.upper()
    p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = accent
    p.font.name = "Helvetica"
    tb = s.shapes.add_textbox(Inches(x + 0.35), Inches(y + 0.5),
                                Inches(w - 0.5), Inches(0.6))
    p = tb.text_frame.paragraphs[0]; p.text = title
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = HEAD
    p.font.name = "Helvetica"
    tb = s.shapes.add_textbox(Inches(x + 0.35), Inches(y + 1.2),
                                Inches(w - 0.5), Inches(h - 1.3))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(body):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  " + line
        p.font.size = Pt(13); p.font.color.rgb = INK; p.font.name = "Helvetica"
        p.space_after = Pt(6)

big_card(s, 0.5, 1.8, 6.1, 4.8,
         "May 2026", "GCB submission",
         ["Pick ONE fire model.",
          "Stay on GFED4 + Opt 1 baseline (≈ Model C, possibly retuned).",
          "Goal: ship one model into Global Carbon Budget on time.",
          "Paper-grade rigor not required for this artifact."],
         HEAD)

big_card(s, 6.7, 1.8, 6.1, 4.8,
         "End of summer 2026", "Paper",
         ["Build the full 6-row table.",
          "Calibrate each model on its own dataset; cross-score on the other.",
          "Report which input variables each optimizer keeps (variable selection).",
          "Argue the structural-change finding.",
          "Title direction: \"Advances in remote sensing alter the structural form of fire dynamics.\""],
         ACCENT)

# ============================================================
# Slide 6 — Data hand-off from Lei
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Data hand-off (Lei, 2026-05-04)",
             "global_baseline_modelC_inputs_1997-2016.nc  ·  one file, all Model C inputs")

# Left: file box
file_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.5), Inches(1.7), Inches(4.5), Inches(5.4))
file_box.fill.solid(); file_box.fill.fore_color.rgb = SOFT
file_box.line.color.rgb = RULE; file_box.line.width = Pt(0.75)
tf = file_box.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.18); tf.margin_top = Inches(0.15)
p = tf.paragraphs[0]
p.text = "global_baseline_modelC_inputs_1997-2016.nc"
p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = HEAD; p.font.name = "Helvetica"

specs = [
    "",
    "Source:  ED stock-fire global run (GCB2025 S3 drivers)",
    "Grid:  0.5° global,  360 × 720",
    "Time:  240 monthly steps, 1997-01..2016-12 (no-leap)",
    "Format:  NetCDF4, float32, zlib-4",
    "",
    "VARIABLES (all time, lat, lon):",
    "  D_bar             dryness index  (mm)",
    "  T_air             air temp  (degC)",
    "  P_ann             annual precip  (mm yr-1)",
    "  P_month           monthly precip  (mm month-1, /12 applied)",
    "  GPP_month_ntrl    natural GPP",
    "  GPP_month_scnd    secondary GPP",
    "  GPP_month_past    pasture GPP",
    "  area_frac_*       landuse fractions",
    "  area_burned       ED stock-fire output",
    "                    (annual broadcast to 12 months)",
    "",
    "Pre-applied:  sentinel <= -9990 -> NaN;  P_month / 12",
]
for line in specs:
    p = tf.add_paragraph(); p.text = line
    p.font.size = Pt(10); p.font.color.rgb = INK; p.font.name = "Helvetica"

# Right: pipeline as numbered steps
steps = [
    ("1", "Plug in",         "Wire the NC into reproduce_modelC.py — one canonical input source, replaces the per-variable .npy stack."),
    ("2", "Reproduce",       "Hit Lei's baseline:  loss ~ 0.881  ·  Pearson r ~ 0.63\n10k-trial Optuna refit on 1997–2010.  Sanity gate."),
    ("3", "Diagnose",        "Sahel + Congo undershoot (1.7 %/yr predicted vs 3.1 %/yr GFED).\nTest:  sigmoid saturation · suppression at high D_bar · GPP-hump cap."),
    ("4", "Region-aware fit","Per-biome / per-continent params.  Compare against the global fit on the held-out 2011–2016."),
    ("5", "Defend with maps","Lei's rule:  no metric-only wins.  Every reported gain comes with a map showing WHERE and WHY."),
]
y = 1.7
for num, title, body in steps:
    # number badge
    badge = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                Inches(5.2), Inches(y), Inches(0.5), Inches(0.5))
    badge.fill.solid(); badge.fill.fore_color.rgb = ACCENT
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]; p.text = num; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = WHT
    p.font.name = "Helvetica"

    # text
    tb = s.shapes.add_textbox(Inches(5.85), Inches(y - 0.05),
                                Inches(7.0), Inches(1.05))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = HEAD
    p.font.name = "Helvetica"
    p2 = tf.add_paragraph(); p2.text = body
    p2.font.size = Pt(10); p2.font.color.rgb = SUBT; p2.font.name = "Helvetica"
    y += 1.07

# ============================================================
# Slide 7 — Train / test split
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Workflow", "Train / test split  ·  no leakage")

# Timeline bar
y0 = 3.1
total_w = 12.0
x0 = 0.7
n_years = 20  # 1997..2016
year_w = total_w / n_years

# Train block
train = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                            Inches(x0), Inches(y0),
                            Inches(year_w * 14), Inches(0.7))
train.fill.solid(); train.fill.fore_color.rgb = HEAD; train.line.fill.background()
tf = train.text_frame
p = tf.paragraphs[0]; p.text = "TRAIN  ·  1997 — 2010  (14 yrs)"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = WHT; p.font.name = "Helvetica"

# Test block
test = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                           Inches(x0 + year_w * 14), Inches(y0),
                           Inches(year_w * 6), Inches(0.7))
test.fill.solid(); test.fill.fore_color.rgb = ACCENT; test.line.fill.background()
tf = test.text_frame
p = tf.paragraphs[0]; p.text = "HOLD-OUT  ·  2011 — 2016"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = WHT; p.font.name = "Helvetica"

# Year tick labels
for i, yr in enumerate(range(1997, 2017)):
    tb = s.shapes.add_textbox(Inches(x0 + i * year_w - 0.15),
                                Inches(y0 + 0.8),
                                Inches(year_w + 0.3), Inches(0.3))
    p = tb.text_frame.paragraphs[0]; p.text = str(yr)
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(9); p.font.color.rgb = SUBT; p.font.name = "Helvetica"

# Rules below
rules = [
    "Optimizer sees only TRAIN years.  Hyperparameters and stopping criteria also tuned only on TRAIN.",
    "All reported metrics + maps are computed on the HOLD-OUT.  No re-tuning if hold-out looks bad.",
    "Year-by-year matching, not climatology — preserves interannual variability skill.",
]
y = 4.6
for r in rules:
    bullet = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(0.85), Inches(y + 0.08),
                                 Inches(0.12), Inches(0.12))
    bullet.fill.solid(); bullet.fill.fore_color.rgb = ACCENT
    bullet.line.fill.background()
    tb = s.shapes.add_textbox(Inches(1.1), Inches(y), Inches(11.5), Inches(0.5))
    p = tb.text_frame.paragraphs[0]; p.text = r
    p.font.size = Pt(13); p.font.color.rgb = INK; p.font.name = "Helvetica"
    y += 0.55

# ============================================================
# Slide 8 — Vocabulary: predictor / mechanism / parameter
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Vocabulary",
             "Predictor  ·  Mechanism  ·  Parameter   — three things that get confused")

# 3 columns
def vocab_card(s, x, y, w, h, name, sub, body, accent):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                              Inches(w), Inches(0.7))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
    tf = bar.text_frame
    p = tf.paragraphs[0]; p.text = name; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = WHT; p.font.name = "Helvetica"

    body_box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(x), Inches(y + 0.7),
                                    Inches(w), Inches(h - 0.7))
    body_box.fill.solid(); body_box.fill.fore_color.rgb = SOFT
    body_box.line.color.rgb = RULE; body_box.line.width = Pt(0.5)
    tf = body_box.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.18)
    p = tf.paragraphs[0]; p.text = sub
    p.font.size = Pt(11); p.font.italic = True; p.font.color.rgb = SUBT; p.font.name = "Helvetica"
    p.space_after = Pt(8)
    for line in body:
        p = tf.add_paragraph(); p.text = line
        p.font.size = Pt(12); p.font.color.rgb = INK; p.font.name = "Helvetica"
        p.space_after = Pt(4)

vocab_card(s, 0.5, 1.7, 4.0, 4.6,
           "PREDICTOR",
           "an INPUT to the formula",
           ["• Comes from data",
            "• Different value at every cell + month",
            "• Model C uses 5:",
            "    D_bar  ·  T_air",
            "    P_ann  ·  P_month  ·  GPP",
            "",
            "Recipe analogy: ingredients"],
           HEAD)

vocab_card(s, 4.7, 1.7, 4.0, 4.6,
           "MECHANISM",
           "a PIECE of the formula",
           ["• What we DO with a predictor",
            "• A switch / hump / dampener",
            "• Model C has 6 mechanisms",
            "    (using the 5 predictors)",
            "",
            "Recipe analogy: a step in the recipe",
            "(\"whisk the eggs\")"],
           ACCENT)

vocab_card(s, 8.9, 1.7, 4.0, 4.6,
           "PARAMETER",
           "a KNOB inside a mechanism",
           ["• A number we tune",
            "• Same value used everywhere",
            "• Model C has 12 parameters",
            "    (~2 per mechanism)",
            "• What lives in params.json",
            "",
            "Recipe analogy: amounts (\"2 cups\")"],
           HEAD)

# Footer relationship line
tb = s.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7))
p = tb.text_frame.paragraphs[0]
p.text = "Relationship:    predictors flow INTO mechanisms;   mechanisms are tuned BY parameters."
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = HEAD; p.font.name = "Helvetica"

# ============================================================
# Slide 9 — Anatomy of Model C  (the 6 mechanisms × 12 parameters)
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Anatomy",
             "Model C  =  6 mechanisms  ·  5 predictors  ·  12 parameters")

n_rows = 1 + 6
n_cols = 4
tbl_shape = s.shapes.add_table(n_rows, n_cols,
                                Inches(0.5), Inches(1.7),
                                Inches(12.33), Inches(4.6))
tbl = tbl_shape.table
widths = [0.5, 4.5, 2.0, 5.33]
for i, w in enumerate(widths):
    tbl.columns[i].width = Inches(w)
tbl.rows[0].height = Inches(0.55)
for r in range(1, n_rows):
    tbl.rows[r].height = Inches(0.7)

for c, h in enumerate(["#", "Mechanism — what it asks", "Predictor(s)", "Parameter(s) — knobs"]):
    cell(tbl, 0, c, h, bold=True, fg=WHT, bg=HEAD,
         align=PP_ALIGN.CENTER, size=12)

mechs = [
    ("1", "Dry enough to burn?  (sigmoid)",        "D_bar",   "k1, D_low"),
    ("2", "Too dry to burn?  (reverse sigmoid)",   "D_bar",   "k2, D_high"),
    ("3", "Enough annual rain to grow stuff?",     "P_ann",   "P_half"),
    ("4", "Raining right now?  (dampener)",        "P_month", "pre_dampen_half"),
    ("5", "GPP hump  (not too sparse, not too dense)", "GPP",  "gpp_af, gpp_b, gpp_d"),
    ("6", "Warm enough to ignite?  (sigmoid)",     "T_air",   "ign_k, ign_c"),
]
for i, row in enumerate(mechs):
    r = i + 1
    bg = WHT if i % 2 == 0 else SOFT
    cell(tbl, r, 0, row[0], bold=True, fg=ACCENT, bg=bg,
         align=PP_ALIGN.CENTER, size=14)
    cell(tbl, r, 1, row[1], bg=bg, size=12)
    cell(tbl, r, 2, row[2], bg=bg, fg=HEAD, bold=True,
         align=PP_ALIGN.CENTER, size=12)
    cell(tbl, r, 3, row[3], bg=bg, fg=SUBT, size=11)

# Plus one extra param outside any mechanism — the global exponent
tb = s.shapes.add_textbox(Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.4))
p = tb.text_frame.paragraphs[0]
p.text = "Plus one global \"sharpness dial\":  fire_exp  — exponent on the whole formula  (12 params total)."
p.font.size = Pt(11); p.font.italic = True; p.font.color.rgb = SUBT; p.font.name = "Helvetica"

tb = s.shapes.add_textbox(Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.5))
p = tb.text_frame.paragraphs[0]
p.text = "fire  =  (mechanism 1) × (m. 2) × (m. 3) × (m. 4) × (m. 5) × (m. 6),  all raised to fire_exp"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = HEAD; p.font.name = "Helvetica"

# ============================================================
# Slide 10 — The 2-family view  (and what's missing)
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Coarser view",
             "Zoom out:  6 math pieces   →   2 real-world families   +   1 missing")

def family_box(s, x, y, w, h, title, mechs_list, status, fill, fg=WHT):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = fill
    box.line.color.rgb = fill
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.18)
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = fg; p.font.name = "Helvetica"
    p2 = tf.add_paragraph(); p2.text = status
    p2.font.size = Pt(11); p2.font.italic = True; p2.font.color.rgb = fg; p2.font.name = "Helvetica"
    p2.space_after = Pt(8)
    for m in mechs_list:
        p = tf.add_paragraph(); p.text = "  " + m
        p.font.size = Pt(12); p.font.color.rgb = fg; p.font.name = "Helvetica"
        p.space_after = Pt(2)

family_box(s, 0.5, 1.7, 4.0, 4.5,
           "Climate / environment",
           ["Mech 1 — dry enough?",
            "Mech 2 — too dry?",
            "Mech 3 — annual rain?",
            "Mech 4 — raining now?",
            "Mech 6 — warm enough?"],
           "PRESENT  ✓",
           HEAD)

family_box(s, 4.7, 1.7, 4.0, 4.5,
           "Fuel / vegetation",
           ["Mech 5 — GPP hump",
            "(GPP is just a proxy for biomass)"],
           "PRESENT  ✓  (partial)",
           ACCENT)

family_box(s, 8.9, 1.7, 4.0, 4.5,
           "Ignition source",
           ["• Lightning",
            "• Humans / roads / cropland burning",
            "",
            "Model C has none of these.",
            "ign_c, ign_k are just \"warm enough?\"",
            "— that's still climate, not ignition."],
           "MISSING  ✗",
           RGBColor(0x55, 0x55, 0x55))

# Footer
tb = s.shapes.add_textbox(Inches(0.5), Inches(6.4), Inches(12.3), Inches(1.0))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Classic fire science:  fire needs  Environment × Fuel × Ignition.  Model C has 2 of the 3."
p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = HEAD; p.font.name = "Helvetica"
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "Hurtt's expectation:  GFED5 + stricter optimization criteria will FORCE the optimizer to bring ignition in."
p2.font.size = Pt(11); p2.font.color.rgb = SUBT; p2.font.name = "Helvetica"
p2.alignment = PP_ALIGN.CENTER

# ============================================================
# Slide 11 — GPP vs biomass  (and the upgrade path)
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Fuel proxies",
             "Why GPP alone isn't enough  —  what it would take to add real fuel")

# 2 columns top: GPP vs Biomass
def proxy_card(s, x, y, w, h, title, what_it_is, captures, misses, fill):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = SOFT
    box.line.color.rgb = RULE; box.line.width = Pt(0.5)
    # Header bar
    head = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(0.6))
    head.fill.solid(); head.fill.fore_color.rgb = fill; head.line.fill.background()
    tf = head.text_frame
    p = tf.paragraphs[0]; p.text = title; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = WHT; p.font.name = "Helvetica"

    tb = s.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.75),
                                Inches(w - 0.36), Inches(h - 0.85))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = what_it_is
    p.font.size = Pt(11); p.font.italic = True; p.font.color.rgb = SUBT; p.font.name = "Helvetica"
    p.space_after = Pt(8)
    p = tf.add_paragraph(); p.text = "Captures well:"
    p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = HEAD; p.font.name = "Helvetica"
    for line in captures:
        p = tf.add_paragraph(); p.text = "  ✓  " + line
        p.font.size = Pt(11); p.font.color.rgb = INK; p.font.name = "Helvetica"
    p = tf.add_paragraph(); p.text = ""
    p = tf.add_paragraph(); p.text = "Misses:"
    p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = ACCENT; p.font.name = "Helvetica"
    for line in misses:
        p = tf.add_paragraph(); p.text = "  ✗  " + line
        p.font.size = Pt(11); p.font.color.rgb = INK; p.font.name = "Helvetica"

proxy_card(s, 0.5, 1.7, 6.1, 3.5,
           "GPP   (what we use now)",
           "Rate — kg/m² per YEAR.  How fast plants are producing.",
           ["Grass fires (high productivity, fast turnover)",
            "Productivity in tropical savanna",
            "Cheap — already in our data"],
           ["Doesn't measure standing fuel",
            "Misses dead litter / dry wood",
            "Treats rainforest and savanna as similar"],
           ACCENT)

proxy_card(s, 6.7, 1.7, 6.1, 3.5,
           "Biomass   (what we should add)",
           "Stock — kg/m².  How much plant material is sitting there.",
           ["Forest fires (centuries of accumulated wood)",
            "Boreal / temperate forest dynamics",
            "Direct measure of what burns"],
           ["Doesn't tell flammable vs non-flammable",
            "Needs HEIGHT to know what's reachable",
            "Not yet in Lei's data dump"],
           HEAD)

# Bottom: upgrade path
upbox = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.7))
upbox.fill.solid(); upbox.fill.fore_color.rgb = HEAD
upbox.line.color.rgb = HEAD
tf = upbox.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.18)
p = tf.paragraphs[0]; p.text = "Path to add fuel properly"
p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = WHT; p.font.name = "Helvetica"
p.space_after = Pt(8)
for line in [
    "1.  Ask Lei to include aboveground biomass per landuse in the next data dump (it's already in ED).",
    "2.  Add a 6th predictor (biomass) and a new mechanism — \"is there enough fuel to burn?\"",
    "3.  Optionally add canopy height as a 7th predictor — gates how much biomass is ground-reachable.",
    "Net change:  5 → 7 predictors  ·  6 → 8 mechanisms  ·  12 → ~16 parameters.",
]:
    p = tf.add_paragraph(); p.text = line
    p.font.size = Pt(12); p.font.color.rgb = WHT; p.font.name = "Helvetica"
    p.space_after = Pt(3)

# ============================================================
# Slide 12 — Section divider + Overarching arc
# ============================================================
s = prs.slides.add_slide(BLANK)
# Big section banner
banner = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0), Inches(0), Inches(13.333), Inches(0.7))
banner.fill.solid(); banner.fill.fore_color.rgb = ACCENT; banner.line.fill.background()
tf = banner.text_frame
p = tf.paragraphs[0]; p.text = "DISSERTATION FRAMING  ·  Hurtt meeting 2026-05-05"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = WHT; p.font.name = "Helvetica"

# Overarching question
tb = s.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.5))
p = tb.text_frame.paragraphs[0]; p.text = "OVERARCHING QUESTION"
p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = ACCENT; p.font.name = "Helvetica"

oq_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(0.5), Inches(1.45), Inches(12.3), Inches(0.95))
oq_box.fill.solid(); oq_box.fill.fore_color.rgb = HEAD
oq_box.line.color.rgb = HEAD
tf = oq_box.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.18)
p = tf.paragraphs[0]
p.text = "What is the role of remote sensing and modeling for fire in the carbon cycle?"
p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = WHT; p.font.name = "Helvetica"
p.alignment = PP_ALIGN.CENTER

# 3 chapter cards
def chap_card(s, x, y, w, h, num, kicker, q_title, q_text, accent):
    # Number badge
    badge = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                Inches(x + 0.2), Inches(y + 0.18),
                                Inches(0.55), Inches(0.55))
    badge.fill.solid(); badge.fill.fore_color.rgb = accent
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]; p.text = num; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = WHT; p.font.name = "Helvetica"

    # Card body
    body = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    body.fill.solid(); body.fill.fore_color.rgb = SOFT
    body.line.color.rgb = RULE; body.line.width = Pt(0.5)
    body.shadow.inherit = False

    # Re-add badge above the body
    badge2 = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(x + 0.2), Inches(y + 0.18),
                                 Inches(0.55), Inches(0.55))
    badge2.fill.solid(); badge2.fill.fore_color.rgb = accent
    badge2.line.fill.background()
    tf = badge2.text_frame
    p = tf.paragraphs[0]; p.text = num; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = WHT; p.font.name = "Helvetica"

    # Kicker
    tb = s.shapes.add_textbox(Inches(x + 0.95), Inches(y + 0.2),
                                Inches(w - 1.1), Inches(0.4))
    p = tb.text_frame.paragraphs[0]; p.text = kicker.upper()
    p.font.size = Pt(10); p.font.bold = True; p.font.color.rgb = accent; p.font.name = "Helvetica"

    # Q title
    tb = s.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.85),
                                Inches(w - 0.4), Inches(0.4))
    p = tb.text_frame.paragraphs[0]; p.text = q_title
    p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = HEAD; p.font.name = "Helvetica"

    # Q text (the actual question)
    tb = s.shapes.add_textbox(Inches(x + 0.2), Inches(y + 1.3),
                                Inches(w - 0.4), Inches(h - 1.4))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = q_text
    p.font.size = Pt(11); p.font.italic = True; p.font.color.rgb = INK; p.font.name = "Helvetica"

chap_card(s, 0.5, 2.7, 4.0, 4.4, "1", "Current paper",
           "Structural form",
           "How do advances in remote sensing alter the structural form of fire dynamics?\n\n"
           "(GFED4 vs GFED5;  3 optimization criteria;  6 candidate models)",
           HEAD)
chap_card(s, 4.7, 2.7, 4.0, 4.4, "2", "Next",
           "Forest fire stratification",
           "How can remote sensing be used to improve the discrimination between "
           "ground fires and canopy fires in forested regions?\n\n"
           "(uses height — ED has it, GEDI lidar has it, most models don't)",
           ACCENT)
chap_card(s, 8.9, 2.7, 4.0, 4.4, "3", "Ending on the future",
           "Predictive capacity",
           "Do advances in remote sensing and modeling convey predictive capacity for "
           "future fires?\n\n"
           "(seasonal-to-decadal anomalies, NOT 100-year scenarios — must be validatable)",
           HEAD)

# Footer linking idea
tb = s.shapes.add_textbox(Inches(0.5), Inches(7.15), Inches(12.3), Inches(0.3))
p = tb.text_frame.paragraphs[0]
p.text = "Throughline:  same fire-model lineage runs through all three  ·  each chapter answers a different question with it"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(11); p.font.italic = True; p.font.color.rgb = SUBT; p.font.name = "Helvetica"

# ============================================================
# Slide 13 — Hurtt's framing rules  (titles, questions, proposal structure)
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Hurtt's rules",
             "How to phrase the question so it survives outside the model community")

# Left: title rules
tb = s.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(6.0), Inches(0.4))
p = tb.text_frame.paragraphs[0]; p.text = "TITLE & QUESTION RULES"
p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = ACCENT; p.font.name = "Helvetica"

rules = [
    ("No proper nouns",
     "The title cannot contain  ED · GFED · GEDI · LIDAR.  Reviewers + program managers don't care about your tools."),
    ("The title IS the question",
     "If the title is well-phrased, the science question is already stated."),
    ("Question simple, answer complex",
     "A good question is short and reads cleanly to a non-expert.  The work behind the answer can be technical."),
    ("Intuitively interesting to a non-expert",
     "Test:  read it to someone outside fire/Earth science.  If they say \"so what?\" — rewrite."),
    ("Avoid \"how would you do it\" wording",
     "That's the methods, one tier down.  Phrase what you're DOING, not which tool you'd use."),
]
y = 2.2
for k, v in rules:
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.5), Inches(y), Inches(6.0), Inches(0.85))
    box.fill.solid(); box.fill.fore_color.rgb = SOFT
    box.line.color.rgb = RULE; box.line.width = Pt(0.5)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.18); tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]; p.text = k
    p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = HEAD; p.font.name = "Helvetica"
    p2 = tf.add_paragraph(); p2.text = v
    p2.font.size = Pt(10); p2.font.color.rgb = SUBT; p2.font.name = "Helvetica"
    y += 0.93

# Right: rule of threes — proposal architecture
tb = s.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.8), Inches(0.4))
p = tb.text_frame.paragraphs[0]; p.text = "PROPOSAL STRUCTURE  ·  RULE OF THREES"
p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = ACCENT; p.font.name = "Helvetica"

# Tree
def tree_node(s, x, y, w, h, label, level):
    fill = HEAD if level == 0 else (ACCENT if level == 1 else RGBColor(0x55,0x55,0x55))
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = fill
    box.line.color.rgb = fill
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.06)
    p = tf.paragraphs[0]; p.text = label; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(10) if level >= 1 else Pt(11)
    p.font.bold = True; p.font.color.rgb = WHT; p.font.name = "Helvetica"

# Root
tree_node(s, 8.4, 2.2, 3.0, 0.5, "Overarching question", 0)
# 3 objectives
xs_obj = [7.1, 9.4, 11.7]
for i, x in enumerate(xs_obj):
    tree_node(s, x - 0.55, 3.1, 1.5, 0.45, f"Objective {i+1}", 1)
# 3 tasks per objective
for x in xs_obj:
    for j in range(3):
        tree_node(s, x - 0.55, 3.85 + j * 0.5, 1.5, 0.4, f"Task", 2)

# Connecting lines (simple)
def connect(x1, y1, x2, y2):
    c = s.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = RULE; c.line.width = Pt(0.75)

# root → obj
for x in xs_obj:
    connect(9.9, 2.7, x, 3.1)
# obj → tasks
for x in xs_obj:
    connect(x, 3.55, x, 3.85)
    connect(x, 3.55, x, 4.35)
    connect(x, 3.55, x, 4.85)

# Footer note
tb = s.shapes.add_textbox(Inches(7.0), Inches(6.4), Inches(5.8), Inches(0.8))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = ("1 question  ·  3 objectives  ·  3 tasks each.")
p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = HEAD; p.font.name = "Helvetica"
p2 = tf.add_paragraph()
p2.text = "Pre-proposal needs the top two tiers tightly written; tasks come at the proposal stage."
p2.font.size = Pt(10); p2.font.color.rgb = SUBT; p2.font.name = "Helvetica"

# ============================================================
# Slide 14 — Chapter 1 deep dive
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Chapter 1  ·  current paper",
             "How do advances in remote sensing alter the structural form of fire dynamics?")

# Question banner
qb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                         Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.7))
qb.fill.solid(); qb.fill.fore_color.rgb = HEAD; qb.line.color.rgb = HEAD
tf = qb.text_frame
p = tf.paragraphs[0]
p.text = "How do advances in remote sensing alter the structural form of fire dynamics?"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = WHT; p.font.name = "Helvetica"

# 3 objective cards
def obj_card(s, x, y, w, h, num, title, body):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = SOFT
    box.line.color.rgb = RULE; box.line.width = Pt(0.5)

    head_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(x), Inches(y), Inches(w), Inches(0.55))
    head_bar.fill.solid(); head_bar.fill.fore_color.rgb = ACCENT
    head_bar.line.fill.background()
    tf = head_bar.text_frame
    p = tf.paragraphs[0]; p.text = f"OBJECTIVE  {num}"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = WHT; p.font.name = "Helvetica"

    tb = s.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.7),
                                Inches(w - 0.36), Inches(0.5))
    p = tb.text_frame.paragraphs[0]; p.text = title
    p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = HEAD; p.font.name = "Helvetica"

    tb = s.shapes.add_textbox(Inches(x + 0.18), Inches(y + 1.25),
                                Inches(w - 0.36), Inches(h - 1.35))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(body):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  " + line
        p.font.size = Pt(11); p.font.color.rgb = INK; p.font.name = "Helvetica"
        p.space_after = Pt(4)

obj_card(s, 0.5, 2.7, 4.0, 4.4, "1",
          "Characterize the advances in remote sensing",
          ["GFED4 vs GFED5 — what changed in the burned-area product",
           "More small fires; ~2× burned area",
           "Why the differences exist (sensors, retrievals, masking)"])

obj_card(s, 4.7, 2.7, 4.0, 4.4, "2",
          "Evaluate influence on model optimization",
          ["Calibrate Model C against GFED4 and against GFED5",
           "Cross 3 optimization criteria  (every-fire / -continent / -fire-type)",
           "= 6 independently-fit models  (the table on slide 2)"])

obj_card(s, 8.9, 2.7, 4.0, 4.4, "3",
          "Attribute changes back to the advances",
          ["Compare which predictors each optimizer KEEPS",
           "Same parameters  →  parametric change.  Different predictors  →  structural change",
           "→ paper's claim:  better remote sensing  →  richer required model structure"])

# ============================================================
# Slide 15 — Chapter 2  ·  forest fire stratification (lidar / height)
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Chapter 2  ·  next",
             "How can remote sensing improve discrimination of ground vs canopy fires in forests?")

qb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                         Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.7))
qb.fill.solid(); qb.fill.fore_color.rgb = ACCENT; qb.line.color.rgb = ACCENT
tf = qb.text_frame
p = tf.paragraphs[0]
p.text = "How can remote sensing be used to improve the discrimination between ground fires and canopy fires in forested regions?"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = WHT; p.font.name = "Helvetica"

# Two columns: why this fits + what we'd do
def info_block(s, x, y, w, h, title, lines, accent):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(x), Inches(y), Inches(0.12), Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
    tb = s.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.05),
                                Inches(w - 0.4), Inches(0.5))
    p = tb.text_frame.paragraphs[0]; p.text = title
    p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = HEAD; p.font.name = "Helvetica"
    tb = s.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.55),
                                Inches(w - 0.4), Inches(h - 0.6))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  " + line
        p.font.size = Pt(11); p.font.color.rgb = INK; p.font.name = "Helvetica"
        p.space_after = Pt(5)

info_block(s, 0.5, 2.7, 6.1, 4.5,
            "Why this fits us specifically",
            ["ED already carries canopy HEIGHT in its state — most DGVMs don't",
             "GEDI satellite lidar measures forest canopy height directly",
             "We have both data sources; few groups can pair them",
             "Hurtt's exact words: \"right in our wheelhouse\""],
            HEAD)
info_block(s, 6.7, 2.7, 6.1, 4.5,
            "What it would do  /  caveats",
            ["Use ED state + GEDI height to predict whether fire is canopy vs ground fire",
             "Restrict to forested regions — at the global scale grasslands swamp the signal",
             "More focused version of Ch 1's structural-form question, applied to forest",
             "Caveat:  scope-limit early — Hurtt warned against trying it globally"],
            ACCENT)

# ============================================================
# Slide 16 — Chapter 3  ·  predictive capacity for future fires
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Chapter 3  ·  ending on the future",
             "Do advances in remote sensing and modeling convey predictive capacity for future fires?")

qb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                         Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.7))
qb.fill.solid(); qb.fill.fore_color.rgb = HEAD; qb.line.color.rgb = HEAD
tf = qb.text_frame
p = tf.paragraphs[0]
p.text = "Do advances in remote sensing and modeling convey predictive capacity for future fires?"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = WHT; p.font.name = "Helvetica"

# Time-scale spectrum bar — what counts and what doesn't
tb = s.shapes.add_textbox(Inches(0.5), Inches(2.6), Inches(12.3), Inches(0.4))
p = tb.text_frame.paragraphs[0]; p.text = "TIME-SCALE LADDER  ·  what's actually predictable"
p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = ACCENT; p.font.name = "Helvetica"

ladders = [
    ("Daily weather",            "skillful out to ~10 days",         GRN := RGBColor(0x1E,0x8A,0x4A)),
    ("Seasonal anomaly  (1–9 mo)", "skillful — used in harvest forecasts (crop yields)", GRN),
    ("Multi-year (decadal)",     "modest skill via ENSO / climate oscillations",         RGBColor(0xE6,0x7E,0x22)),
    ("100-year scenarios",       "intellectually interesting BUT impossible to validate — Hurtt steered us away", RGBColor(0xC0,0x39,0x2B)),
]
y = 3.1
for label, note, color in ladders:
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.5), Inches(y), Inches(0.18), Inches(0.55))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.85), Inches(y + 0.04),
                                Inches(4.0), Inches(0.5))
    p = tb.text_frame.paragraphs[0]; p.text = label
    p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = HEAD; p.font.name = "Helvetica"
    tb = s.shapes.add_textbox(Inches(4.9), Inches(y + 0.07),
                                Inches(8.0), Inches(0.5))
    p = tb.text_frame.paragraphs[0]; p.text = note
    p.font.size = Pt(11); p.font.color.rgb = SUBT; p.font.name = "Helvetica"
    y += 0.65

# Use cases
uc_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(0.5), Inches(5.95), Inches(12.3), Inches(1.3))
uc_box.fill.solid(); uc_box.fill.fore_color.rgb = HEAD; uc_box.line.color.rgb = HEAD
tf = uc_box.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.15)
p = tf.paragraphs[0]; p.text = "Why anyone cares  (Hurtt's harvest analogy)"
p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = WHT; p.font.name = "Helvetica"
p.space_after = Pt(6)
for line in [
    "Harvest forecasters already convert seasonal weather skill into 6–9 month crop-yield predictions.  Why not fire?",
    "If you can warn 6 months ahead that a region will have a high fire season → budget firefighting · stockpile resources · ban campfires · prescribed burns · fuel-thinning.",
]:
    p = tf.add_paragraph(); p.text = line
    p.font.size = Pt(11); p.font.color.rgb = WHT; p.font.name = "Helvetica"
    p.space_after = Pt(3)

# ============================================================
# Slide 17 — Peer feedback on pre-proposal outline (2026-05-07)
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Peer feedback  ·  pre-proposal outline",
             "Three suggestions from a peer expert reviewer  (received 2026-05-07)")

# Status badge helper
def status_badge(s, x, y, w, h, label, fill):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = fill
    box.line.fill.background()
    tf = box.text_frame
    p = tf.paragraphs[0]; p.text = label; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(9); p.font.bold = True; p.font.color.rgb = WHT; p.font.name = "Helvetica"

# Three feedback cards
def fb_card(s, x, y, w, h, num, kicker, body, status, status_color):
    body_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(x), Inches(y), Inches(w), Inches(h))
    body_box.fill.solid(); body_box.fill.fore_color.rgb = SOFT
    body_box.line.color.rgb = RULE; body_box.line.width = Pt(0.5)

    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(x), Inches(y), Inches(w), Inches(0.6))
    bar.fill.solid(); bar.fill.fore_color.rgb = HEAD; bar.line.fill.background()
    tf = bar.text_frame
    p = tf.paragraphs[0]; p.text = f"  #{num}    {kicker}"
    p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = WHT; p.font.name = "Helvetica"

    # status pill in top-right
    status_badge(s, x + w - 1.5, y + 0.13, 1.35, 0.34, status, status_color)

    tb = s.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.75),
                                Inches(w - 0.36), Inches(h - 0.85))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(body):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11); p.font.color.rgb = INK; p.font.name = "Helvetica"
        p.space_after = Pt(5)

fb_card(s, 0.5, 1.7, 4.0, 5.0, "1",
         "Pareto front (Ch 1)",
         ['"Sometimes a model can fit the data well in two different ways. '
          'Showing where the model fails to reconcile both will be your strongest '
          'evidence for structural gaps."',
          "",
          "Action: add Pareto-front analysis to Ch 1 Obj 2/3 at task stage. "
          "Pairs naturally with the planned Metropolis-Hastings ensemble."],
         "PENDING", RGBColor(0xE6, 0x7E, 0x22))

fb_card(s, 4.7, 1.7, 4.0, 5.0, "2",
         "Amazon edge case (Ch 2)",
         ['"The Amazon edge is where height-explicit demographics and fire '
          'intensity collide. Highlighting this specific biome would ground '
          'the global scope in a high-stakes reality."',
          "",
          "Action: name the Amazon edge as a primary test region in Ch 2 "
          "motivation and Objective 3."],
         "PENDING", RGBColor(0xE6, 0x7E, 0x22))

fb_card(s, 8.9, 1.7, 4.0, 5.0, "3",
         "GFED5 first-mover claim",
         ['"GFED5 is very new. Emphasize that your research is the first '
          'to systematically evaluate its structural implications. This '
          'first-mover advantage makes the proposal very compelling."',
          "",
          "Action: added one sentence to Ch 1 Motivation flagging the "
          "first-mover claim."],
         "INCORPORATED", RGBColor(0x1E, 0x8A, 0x4A))

# Footer with reviewer's overall verdict
verd = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.45))
verd.fill.solid(); verd.fill.fore_color.rgb = HEAD; verd.line.color.rgb = HEAD
tf = verd.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.18); tf.margin_top = Inches(0.06)
p = tf.paragraphs[0]
p.text = ('Overall verdict:  "You aren\'t just building a better fire model. '
          'You are using modern data and optimization to figure out why previous '
          'models were structurally limited."')
p.font.size = Pt(11); p.font.italic = True; p.font.color.rgb = WHT; p.font.name = "Helvetica"

# ============================================================
# Slide 18 — Pre-proposal schematic (one-slide concept diagram)
# ============================================================
s = prs.slides.add_slide(BLANK)

# --- Top carbon-cycle frame ---
top_band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0), Inches(0), Inches(13.333), Inches(0.42))
top_band.fill.solid(); top_band.fill.fore_color.rgb = HEAD
top_band.line.fill.background()
tf = top_band.text_frame
p = tf.paragraphs[0]
p.text = "FIRE IN THE CARBON CYCLE  ·  THE THROUGHLINE OF THIS DISSERTATION"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = WHT
p.font.name = "Helvetica"

# --- Overarching question ---
oq = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                         Inches(0.5), Inches(0.6), Inches(12.3), Inches(0.85))
oq.fill.solid(); oq.fill.fore_color.rgb = SOFT
oq.line.color.rgb = ACCENT; oq.line.width = Pt(1.5)
tf = oq.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.05)
p = tf.paragraphs[0]
p.text = "OVERARCHING QUESTION"
p.font.size = Pt(9); p.font.bold = True; p.font.color.rgb = ACCENT
p.font.name = "Helvetica"
p2 = tf.add_paragraph()
p2.text = "What is the role of remote sensing and modeling for fire in the carbon cycle?"
p2.font.size = Pt(17); p2.font.bold = True; p2.font.color.rgb = HEAD
p2.font.name = "Helvetica"

# --- Three chapter cards with arrows between ---
def chap_panel(s, x, y, w, h, num, title, question, objectives):
    # Body
    body = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(x), Inches(y), Inches(w), Inches(h))
    body.fill.solid(); body.fill.fore_color.rgb = SOFT
    body.line.color.rgb = RULE; body.line.width = Pt(0.5)

    # Header
    head_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(x), Inches(y), Inches(w), Inches(0.55))
    head_bar.fill.solid(); head_bar.fill.fore_color.rgb = HEAD
    head_bar.line.fill.background()
    tf = head_bar.text_frame
    p = tf.paragraphs[0]; p.text = f"  CHAPTER {num}    {title}"
    p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = WHT
    p.font.name = "Helvetica"

    # Question
    tb = s.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.7),
                                Inches(w - 0.36), Inches(0.85))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = question
    p.font.size = Pt(10.5); p.font.italic = True; p.font.color.rgb = ACCENT
    p.font.name = "Helvetica"

    # Objectives label
    tb = s.shapes.add_textbox(Inches(x + 0.18), Inches(y + 1.65),
                                Inches(w - 0.36), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.text = "OBJECTIVES"
    p.font.size = Pt(9); p.font.bold = True; p.font.color.rgb = SUBT
    p.font.name = "Helvetica"

    # Objectives
    tb = s.shapes.add_textbox(Inches(x + 0.18), Inches(y + 1.95),
                                Inches(w - 0.36), Inches(h - 2.05))
    tf = tb.text_frame; tf.word_wrap = True
    for i, obj in enumerate(objectives):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{i+1}.  {obj}"
        p.font.size = Pt(10); p.font.color.rgb = INK; p.font.name = "Helvetica"
        p.space_after = Pt(5)

# Layout positions
chap_y = 1.65
chap_h = 3.4
chap_w = 3.7
gaps_x = [0.5, 4.6, 8.7]   # left edges of three cards
arrow_x = [(4.2, 4.55), (8.3, 8.65)]  # arrows between cards

chap_panel(s, gaps_x[0], chap_y, chap_w, chap_h, "1",
            "Structural form of fire dynamics",
            "How do advances in remote sensing alter the structural form of fire dynamics?",
            ["Characterize how the new burned-area data differ from the old.",
             "Determine whether the model needs new predictors or just new parameters.",
             "Attribute model changes back to specific data advances."])

chap_panel(s, gaps_x[1], chap_y, chap_w, chap_h, "2",
            "Vegetation structure & fire severity",
            "How can remote sensing improve discrimination of ground vs canopy fires in forests?",
            ["Identify remote-sensing signatures of ground vs canopy fires.",
             "Evaluate model height against observed canopy structure under fire.",
             "Establish a height-aware fire-typing framework."])

chap_panel(s, gaps_x[2], chap_y, chap_w, chap_h, "3",
            "Predictive capacity for future fires",
            "Do advances in remote sensing and modeling convey predictive capacity for future fires?",
            ["Produce seasonal fire-anomaly forecasts at 1 to 9 month lead.",
             "Quantify forecast skill by region, biome, and lead time.",
             "Show how forecasts inform real fire-management decisions."])

# Arrows between cards
def chain_arrow(s, x1, x2, y_mid):
    arr = s.shapes.add_connector(1, Inches(x1), Inches(y_mid),
                                    Inches(x2), Inches(y_mid))
    arr.line.color.rgb = ACCENT; arr.line.width = Pt(2.5)
    try:
        from pptx.oxml.ns import qn
        from lxml import etree
        ln = arr.line._get_or_add_ln()
        tail = etree.SubElement(ln, qn('a:tailEnd'))
        tail.set('type', 'triangle')
    except Exception:
        pass

mid_y = chap_y + chap_h / 2
for x1, x2 in arrow_x:
    chain_arrow(s, x1, x2, mid_y)

# Mini caption under each arrow
def arr_caption(s, x_center, y, text):
    tb = s.shapes.add_textbox(Inches(x_center - 0.7), Inches(y),
                                Inches(1.4), Inches(0.3))
    p = tb.text_frame.paragraphs[0]; p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(9); p.font.italic = True; p.font.color.rgb = ACCENT
    p.font.name = "Helvetica"

arr_caption(s, 4.375, mid_y - 0.5, "feeds")
arr_caption(s, 8.475, mid_y - 0.5, "feeds")

# --- Foundation band (shared inputs) ---
fnd_y = 5.25
fnd_band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0.5), Inches(fnd_y),
                                Inches(12.3), Inches(0.85))
fnd_band.fill.solid(); fnd_band.fill.fore_color.rgb = SOFT
fnd_band.line.color.rgb = RULE; fnd_band.line.width = Pt(0.5)

# Label on the left
tb = s.shapes.add_textbox(Inches(0.65), Inches(fnd_y + 0.08),
                            Inches(2.0), Inches(0.7))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "FOUNDATION"
p.font.size = Pt(10); p.font.bold = True; p.font.color.rgb = ACCENT
p.font.name = "Helvetica"
p2 = tf.add_paragraph(); p2.text = "(all inputs already accessible)"
p2.font.size = Pt(9); p2.font.italic = True; p2.font.color.rgb = SUBT
p2.font.name = "Helvetica"

# Chips of input data
chips = ["GFED4.1s", "GFED5", "GEDI canopy", "ED model state",
          "CRUJRA climate", "Seasonal forecasts", "Optimization tools"]
chip_x = 2.7
chip_y = fnd_y + 0.22
chip_w = 1.42
chip_h = 0.42
chip_gap = 0.05
for i, c in enumerate(chips):
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(chip_x), Inches(chip_y),
                               Inches(chip_w), Inches(chip_h))
    chip.fill.solid(); chip.fill.fore_color.rgb = HEAD
    chip.line.fill.background()
    tf = chip.text_frame
    p = tf.paragraphs[0]; p.text = c; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(9); p.font.bold = True; p.font.color.rgb = WHT
    p.font.name = "Helvetica"
    chip_x += chip_w + chip_gap

# --- Outputs band ---
out_y = 6.2
out_band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0.5), Inches(out_y),
                                Inches(12.3), Inches(0.85))
out_band.fill.solid(); out_band.fill.fore_color.rgb = HEAD
out_band.line.fill.background()

tb = s.shapes.add_textbox(Inches(0.65), Inches(out_y + 0.08),
                            Inches(2.0), Inches(0.7))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "OUTPUTS"
p.font.size = Pt(10); p.font.bold = True; p.font.color.rgb = WHT
p.font.name = "Helvetica"
p2 = tf.add_paragraph(); p2.text = "(what the dissertation delivers)"
p2.font.size = Pt(9); p2.font.italic = True; p2.font.color.rgb = WHT
p2.font.name = "Helvetica"

outputs = [
    "3 peer-reviewed papers",
    "Integrated, observation-driven, prediction-capable fire model",
    "Open methodological workflow (transferable to other DGVMs)",
    "Operational outputs for fire management, carbon accounting, climate policy",
]
out_x = 2.7
out_w_total = 12.3 - (out_x - 0.5) - 0.15
out_each = (out_w_total - 3 * 0.08) / 4
for i, o in enumerate(outputs):
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(out_x), Inches(out_y + 0.18),
                               Inches(out_each), Inches(0.5))
    chip.fill.solid(); chip.fill.fore_color.rgb = WHT
    chip.line.color.rgb = ACCENT; chip.line.width = Pt(1.0)
    tf = chip.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.06)
    p = tf.paragraphs[0]; p.text = o; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(9); p.font.bold = True; p.font.color.rgb = HEAD
    p.font.name = "Helvetica"
    out_x += out_each + 0.08

# --- Bottom carbon-cycle frame ---
bot_band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0), Inches(7.1),
                                Inches(13.333), Inches(0.4))
bot_band.fill.solid(); bot_band.fill.fore_color.rgb = ACCENT
bot_band.line.fill.background()
tf = bot_band.text_frame
p = tf.paragraphs[0]
p.text = "STRENGTHENS THE REPRESENTATION OF FIRE IN THE GLOBAL CARBON CYCLE"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = WHT
p.font.name = "Helvetica"

prs.save("figures_and_tables.pptx")
print("wrote figures_and_tables.pptx")
