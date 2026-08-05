"""Assemble the GDP -> Model F advisor deck (for the weekly George meeting).

Story = the four beats of the human-factor assignment, GDP only (no C/D/E backstory,
George knows it): plot it -> prove it is real -> prove it is the right factor ->
build it in biome-specifically -> it wins. Speaker notes carry the spoken narrative +
the backup numbers so the deck is presentable as-is.

Figures come from make_advisor_gdp_deck_figs.py (run that first). -> GDP_ModelF_advisor_deck.pptx
"""
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

FIG = Path("figs_gdp_advisor"); OUT = "GDP_ModelF_advisor_deck.pptx"
INK = RGBColor(0x22, 0x33, 0x44); FIRE = RGBColor(0xC1, 0x44, 0x2E)
GREY = RGBColor(0x6b, 0x72, 0x80); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
BLANK = prs.slide_layouts[6]

def tb(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(x, y, w, h); tf = box.text_frame; tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = font
    return box

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def fig_slide(title, subtitle, img, note, tag):
    s = prs.slides.add_slide(BLANK)
    # top accent bar
    bar = s.shapes.add_shape(1, 0, 0, SW, Inches(0.12))
    bar.fill.solid(); bar.fill.fore_color.rgb = FIRE; bar.line.fill.background()
    tb(s, Inches(0.5), Inches(0.22), Inches(11.3), Inches(0.7), title, 26, INK, bold=True)
    tb(s, Inches(0.5), Inches(0.92), Inches(11.3), Inches(0.4), subtitle, 15, GREY)
    # step tag top-right
    tb(s, Inches(11.7), Inches(0.24), Inches(1.5), Inches(0.4), tag, 13, FIRE, bold=True, align=PP_ALIGN.RIGHT)
    # figure, aspect-preserved, fit into the content box
    iw, ih = Image.open(img).size; ar = iw / ih
    box_w, box_h = Inches(12.0), Inches(5.55); bx, by = Inches(0.66), Inches(1.55)
    if box_w / box_h > ar:
        h = box_h; w = int(box_h * ar); x = bx + (box_w - w) // 2; y = by
    else:
        w = box_w; h = int(box_w / ar); x = bx; y = by + (box_h - h) // 2
    s.shapes.add_picture(str(img), x, y, width=w, height=h)
    notes(s, note)
    return s

# ---------------- Slide 1 — title ----------------
s = prs.slides.add_slide(BLANK)
band = s.shapes.add_shape(1, 0, Inches(2.5), SW, Inches(2.5))
band.fill.solid(); band.fill.fore_color.rgb = INK; band.line.fill.background()
tb(s, Inches(0.8), Inches(2.85), Inches(11.7), Inches(1.1),
   "Bringing the human factor into the fire model", 36, WHITE, bold=True)
tb(s, Inches(0.8), Inches(3.95), Inches(11.7), Inches(0.6),
   "A socioeconomic (GDP) fire-suppression term -> Model F", 20, RGBColor(0xE8, 0xC9, 0xC2))
tb(s, Inches(0.8), Inches(5.35), Inches(11.7), Inches(0.5),
   "Richard Owusu-Ansah   |   weekly fire meeting   |   the GDP assignment from 07/23", 14, GREY)
notes(s, "One-line frame before the beats: 'You asked me to test whether human wealth affects "
         "fire and build it into the model. Here is what I did, in four steps -- and it became my "
         "best model, the one now with Lei for the coupled run.' Then go to slide 2. Do NOT relitigate "
         "C/D/E; George knows the backstory. This is only the human-factor step.")

# ---------------- Slide 2 — beat 1 ----------------
fig_slide(
    "The wealth signal is real, not climate in disguise",
    "Fire vs GDP for every country -- and it survives removing climate",
    FIG / "beat1_fire_vs_gdp.png",
    "BEAT 1 -- I built the plot you asked for: fire versus GDP for every country. Fire clearly "
    "falls as countries get richer (left, slope -0.92/decade, r -0.55, 164 countries). Then the "
    "obvious objection: maybe that is just because poor countries happen to be hot dry savanna. So "
    "I removed the climate effect -- temperature, rainfall, vegetation -- and looked again. Wealth "
    "STILL predicts less fire (right, slope -0.70/decade, r -0.47, p = 2e-10). Roughly three "
    "quarters of the raw slope survives. The signal is socioeconomic, not climate wearing an "
    "economics mask.",
    "Beat 1 of 4")

# ---------------- Slide 3 — beat 2 ----------------
fig_slide(
    "GDP is the only human factor that adds anything on its own",
    "Population and land use bring no independent signal once GDP is in",
    FIG / "beat2_controls.png",
    "BEAT 2 -- I checked it was not just GDP getting lucky, and that I was not missing a better "
    "human variable. Population density: no independent signal beyond climate+GDP (F=0.1, p=0.93) "
    "-- the cloud is flat. Land use (pasture/managed land): real on its own, but redundant once GDP "
    "is in (F=1.9, p=0.13), and the model already ingests land use through GPP. So GDP is the one "
    "human factor with independent skill. That is why the model has a GDP term and not a population "
    "or land-use term -- it is a deliberate choice, backed by the negative controls.",
    "Beat 2 of 4")

# ---------------- Slide 4 — beat 3 ----------------
fig_slide(
    "I built it in biome-specifically, not as one global knob",
    "Strong where people manage fire (Africa), near zero where it would erase real fire (Asia)",
    FIG / "beat3_gamma_map.png",
    "BEAT 3 -- Then I built it into the model, biome-specific like you said, not one global knob. "
    "The term is strong in African savanna (gamma 1.6), where people actively manage and suppress "
    "fire, and near zero in poor wet monsoon Asia (0.1) and Australia (0), where cranking the same "
    "knob would wrongly erase real fire. The field is smooth -- Gaussian-blended, no hard borders -- "
    "so it drops straight into Lei's coupled run without seams propagating into the vegetation. "
    "HONEST NOTE if asked: these regional strengths are calibrated to GFED5, not derived from first "
    "principles -- normal for this kind of model, but say it before he does.",
    "Beat 3 of 4")

# ---------------- Slide 5 — beat 4 ----------------
fig_slide(
    "The payoff: Model F reproduces the regional pattern",
    "Right region by region, not a global total hiding compensating errors",
    FIG / "beat4_regional.png",
    "BEAT 4 -- The result is Model F. Adding the human term raised the score AND -- the important "
    "part -- the model now matches the regional pattern, not just the global total (0.98x global). "
    "A model can get the global number right while being wrong everywhere, errors cancelling; this "
    "is not that. HONEST NOTE: boreal (0.52x) and Australia (0.37x) still run low. That is a "
    "FUEL/biomass limitation in the base model, not the human term -- boreal is fuel-limited and a "
    "future base-fuel fix, not a GDP fix. Name it as the known open item.",
    "Beat 4 of 4")

# ---------------- Slide 6 — bottom line ----------------
s = prs.slides.add_slide(BLANK)
bar = s.shapes.add_shape(1, 0, 0, SW, Inches(0.12))
bar.fill.solid(); bar.fill.fore_color.rgb = FIRE; bar.line.fill.background()
tb(s, Inches(0.5), Inches(0.35), Inches(12), Inches(0.8), "Bottom line", 30, INK, bold=True)
bl = ("Model F = the fire model plus a socioeconomic (GDP) suppression term.\n\n"
      "-  The wealth signal is real: it survives removing climate (-0.70/decade, p = 2e-10).\n"
      "-  GDP is the right factor: population (p = 0.93) and land use (p = 0.13) add nothing on their own.\n"
      "-  Built biome-specific and smooth -- coupling-ready, no seams.\n"
      "-  Best model to date on BOTH scored fire rows: burned area 0.679, emissions 0.667.\n"
      "-  Already handed to Lei for the coupled GCB / TRENDY run.\n\n"
      "Open item: boreal + Australia run low -- a base fuel/biomass limit, not the human term.")
tb(s, Inches(0.7), Inches(1.4), Inches(12), Inches(5.5), bl, 18, INK)
notes(s, "Close: 'So the human factor is real, it is the right one, it is built in cleanly, and it "
         "gave my best model on both burned area and emissions -- now with Lei for the coupled run. "
         "The one open item is boreal and Australia fuel, which is a base-model fix, not a human-term "
         "fix.' Controlled GDP-only score jump, if he asks how much GDP alone bought: 0.6547 -> 0.6603.")

prs.save(OUT)
print("wrote", OUT, "--", len(prs.slides._sldIdLst), "slides")
