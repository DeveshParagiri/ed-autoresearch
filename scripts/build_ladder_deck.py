"""Advisor talking-aid deck for the C -> D -> E model ladder.

Short glanceable slides + speaker notes with the plain-language script.
Palette matches build_deck.py. Output: ModelCDE_advisor_deck.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

INK    = RGBColor(0x1A, 0x1A, 0x1A)
SUBT   = RGBColor(0x6B, 0x72, 0x80)
RULE   = RGBColor(0xD0, 0xD5, 0xDB)
HEAD   = RGBColor(0x11, 0x2D, 0x4E)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)
GREEN  = RGBColor(0x1E, 0x7A, 0x5A)
SOFT   = RGBColor(0xF4, 0xF6, 0xF8)
WHT    = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]
W = 13.333


def tb(slide, x, y, w, h, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
       italic=False, space_after=6):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln; p.alignment = align
        p.space_after = Pt(space_after)
        for r in p.runs:
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.italic = italic
            r.font.name = "Calibri"
    return box


def slide_head(title, kicker=None):
    s = prs.slides.add_slide(BLANK)
    if kicker:
        tb(s, 0.6, 0.35, 12, 0.4, kicker.upper(), 13, SUBT, bold=True)
    tb(s, 0.6, 0.72, 12.2, 0.9, title, 32, HEAD, bold=True)
    ln = s.shapes.add_shape(1, Inches(0.6), Inches(1.62), Inches(12.1), Inches(0.03))
    ln.fill.solid(); ln.fill.fore_color.rgb = RULE; ln.line.fill.background()
    ln.shadow.inherit = False
    return s


def bullets(s, items, top=2.0, size=19, gap=0.62, x=0.9, w=11.8):
    for i, it in enumerate(items):
        bold = it.startswith("*")
        txt = it[1:] if bold else it
        col = ACCENT if bold else INK
        tb(s, x, top + i * gap, w, 0.55, txt, size, col, bold=bold)


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


# ---------------------------------------------------------------- 1 title
s = prs.slides.add_slide(BLANK)
bar = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(W), Inches(2.6))
bar.fill.solid(); bar.fill.fore_color.rgb = HEAD; bar.line.fill.background(); bar.shadow.inherit = False
tb(s, 0.8, 0.85, 12, 0.8, "The model ladder", 40, WHT, bold=True)
tb(s, 0.8, 1.68, 12, 0.6, "Model C  ->  Model D  ->  Model E", 24, RGBColor(0xAF, 0xC6, 0xDB))
tb(s, 0.8, 3.1, 12, 0.6, "Each step changes ONE lever", 26, INK, bold=True)
bullets(s, ["C  ->  D    changed the goodness-of-fit CRITERION  (what we optimize)",
            "D  ->  E    changed the FORM  (the equation itself)"], top=3.9, size=21, gap=0.62)
tb(s, 0.9, 5.5, 12, 0.6,
   "Result: the FORM was the binding constraint, not the criterion.", 22, ACCENT, bold=True)
tb(s, 0.8, 6.7, 12, 0.4, "Richard Owusu-Ansah   |   ED fire submodule   |   scored against GFED5 via ILAMB",
   13, SUBT)
notes(s, "Open with: each version differs by ONE LEVER, form or criterion. "
         "Say 'lever' not 'change' - E's form change is a bundle of four edits.")

# ---------------------------------------------------------------- 2 setup
s = slide_head("What the ladder is testing", "the question")
bullets(s, [
    "Fire skill in the model was poor. Two possible explanations:",
    "*1.  We were optimizing the wrong target  (the fit criterion)",
    "*2.  The equation itself could not represent the physics  (the form)",
    "",
    "The ladder separates them. Change one lever at a time, see which one moves the result.",
], top=2.1, gap=0.6)
tb(s, 0.9, 5.5, 11.8, 0.8,
   "Your 1:1 bar is the yardstick throughout:  model vs GFED5 per-cell burned fraction,\n"
   "and the slope of that scatter should be 1.", 19, HEAD, bold=True)
notes(s, "Frame it as a controlled experiment. Two candidate explanations, one lever at a time. "
         "The yardstick is George's own 1:1 line.")

# ---------------------------------------------------------------- 3 model C
s = slide_head("Model C  -  the baseline", "starting point")
bullets(s, [
    "Three-mechanism closed-form equation  (dryness, precipitation, fuel/GPP, ignition)",
    "A product of bounded [0,1] terms, raised to a power, giving an annual fire rate",
    "*Tuned to maximise ILAMB Overall  -  the aggregate community score",
    "",
    "ILAMB Overall  0.649        per-cell 1:1 slope  0.40        magnitude 1.26x GFED5",
], top=2.1, gap=0.62)
tb(s, 0.9, 5.6, 11.8, 0.7,
   "Good aggregate score, but the 1:1 slope sits at 0.40 when it should be 1.", 20, ACCENT, bold=True)
notes(s, "C is the reference. Physical, interpretable equation, tuned to the standard ILAMB score. "
         "It scores respectably but fails the 1:1 bar.")

# ---------------------------------------------------------------- 4 model D
s = slide_head("Model D  -  changed the CRITERION", "lever 1  |  same equation as C")
bullets(s, [
    "*Same equation. Only the optimisation target changed.",
    "New target = spatial Taylor skill, computed only on cells that actually burn",
    "   -  does the fire go in the right PLACES  (correlation r)",
    "   -  are the busy places as busy as reality  (dynamic range, sigma)",
], top=2.05, gap=0.55)
box = s.shapes.add_shape(1, Inches(0.9), Inches(4.35), Inches(11.8), Inches(1.25))
box.fill.solid(); box.fill.fore_color.rgb = SOFT; box.line.color.rgb = RULE; box.shadow.inherit = False
tb(s, 1.15, 4.5, 11.3, 1.0,
   "Why:   ILAMB Overall  =  (Bias + 2xRMSE + Seasonal + Spatial) / 5\n"
   "The spatial pattern was only ONE FIFTH of what we optimised, and RMSE counted double.",
   18, HEAD, bold=True)
tb(s, 0.9, 5.95, 11.8, 0.6, "ILAMB Overall  0.641        1:1 slope  0.42   (barely moved)", 19, INK)
tb(s, 0.9, 6.5, 11.8, 0.5, "Came straight out of your June 9 fire meeting  -  "
   "\"we are missing a very strong spatial pattern\"", 17, SUBT, italic=True)
notes(s, "KEY LINE: we were tuning to a number that barely cared about the thing you were judging us on. "
         "Spatial was 1/5 of the target, RMSE counted double. "
         "If asked 'did you invent the metric' - no, it is the Taylor skill score, the same formula ILAMB "
         "uses for its own Spatial Distribution Score. We computed it only where fire happens and made it "
         "the target instead of one fifth of an average. "
         "If asked why ILAMB went DOWN - because we stopped optimising that number. That is the finding.")

# ---------------------------------------------------------------- 5 model E
s = slide_head("Model E  -  changed the FORM", "lever 2  |  criterion held at D's")
bullets(s, [
    "*Four physical edits to the equation, criterion unchanged from D:",
    "1.  Let the fire rate exceed 1 / yr   (fuel-scaled amplitude)",
    "2.  Fixed the monthly conversion so dry-season fire concentrates",
    "3.  Tropical closed-canopy suppression  (kills false Amazon / Congo fire)",
    "4.  Per-continent parameters",
], top=2.05, gap=0.55)
tb(s, 0.9, 5.15, 11.8, 0.6,
   "ILAMB Overall  0.665        1:1 slope  0.66        r  0.71        magnitude 1.03x", 20, INK, bold=True)
tb(s, 0.9, 5.9, 11.8, 0.9,
   "From your meeting: item A said the rate must be able to exceed 1;  item C said\n"
   "\"different fires need different models for different continents\".", 17, SUBT, italic=True)
notes(s, "Say BUNDLE openly - four edits, one lever. "
         "Edits 1 and 4 are George's own meeting items A and C.")

# ---------------------------------------------------------------- 6 why C could not win
s = slide_head("Why Model C could never reach GFED5", "the structural argument")
bullets(s, [
    "C is a product of [0,1] terms, so the annual rate is hard-capped at 1",
    "The old monthly conversion spread one annual fraction evenly, capping any month at 1/12",
], top=2.1, gap=0.6)
box = s.shapes.add_shape(1, Inches(0.9), Inches(3.35), Inches(11.8), Inches(1.15))
box.fill.solid(); box.fill.fore_color.rgb = SOFT; box.line.color.rgb = RULE; box.shadow.inherit = False
tb(s, 1.15, 3.55, 11.3, 0.9,
   "Chain it through:  per-cell burned fraction caps at 0.039        GFED5 reaches 0.104",
   21, ACCENT, bold=True)
bullets(s, [
    "*No parameter set and no criterion could fix that. It is an architecture problem.",
    "And the correlation r was pinned near 0.5 because continents need OPPOSITE fixes:",
    "   Africa under-burning  |  Amazon over-burning 2x  |  boreal under-burning 4x  |  Europe anti-correlated",
], top=4.75, gap=0.55)
notes(s, "This is the strongest single argument for E. C was structurally incapable, not badly tuned. "
         "Then the r problem: one global formula cannot do opposite things in different places.")

# ---------------------------------------------------------------- 7 africa find
s = slide_head("The find inside Africa", "why the fuel term matters")
bullets(s, [
    "In African savanna, GFED5 fire RISES with productivity  (fuel-limited system)",
    "*Model C had it backwards  -  correlation -0.29",
    "The global GPP hump was suppressing exactly the productive cells that burn most",
    "",
    "*Adding the fuel term took Africa's correlation from 0.47 to 0.66",
], top=2.1, gap=0.62)
tb(s, 0.9, 5.7, 11.8, 0.7,
   "Africa is roughly 60% of global burned area, so fixing its pattern moves everything.",
   20, HEAD, bold=True)
notes(s, "A concrete, physical discovery rather than a tuning story. The model had the fuel "
         "relationship inverted in the single most important fire region.")

# ---------------------------------------------------------------- 8 results table
s = slide_head("The result", "one lever at a time")
rows, cols = 5, 5
tbl = s.shapes.add_table(rows, cols, Inches(1.1), Inches(2.15), Inches(11.1), Inches(2.9)).table
for i, wdt in enumerate([3.1, 2.0, 2.0, 2.0, 2.0]):
    tbl.columns[i].width = Inches(wdt)
hdr = ["", "Model C", "Model D", "Model E", "GFED5"]
data = [["what changed", "baseline", "CRITERION", "FORM", "reference"],
        ["1:1 slope", "0.40", "0.42", "0.66", "1.00"],
        ["correlation r", "~0.50", "~0.50", "0.71", "1.00"],
        ["ILAMB Overall", "0.649", "0.641", "0.665", "-"]]
for j, h in enumerate(hdr):
    c = tbl.cell(0, j); c.text = h
    for p in c.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.size = Pt(17); r.font.bold = True; r.font.color.rgb = WHT
    c.fill.solid(); c.fill.fore_color.rgb = HEAD
for i, row in enumerate(data, start=1):
    for j, v in enumerate(row):
        c = tbl.cell(i, j); c.text = v
        for p in c.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = Pt(16)
                r.font.bold = (j == 3)
                r.font.color.rgb = ACCENT if (j == 3 and i >= 1) else INK
        c.fill.solid(); c.fill.fore_color.rgb = SOFT if i % 2 else WHT
tb(s, 0.9, 5.35, 11.8, 1.2,
   "Criterion alone moved the slope 0.40 -> 0.42.   Changing the form moved it to 0.66.",
   22, ACCENT, bold=True)
tb(s, 0.9, 6.15, 11.8, 0.7,
   "=>  The FORM was the binding constraint. Model D is the control that proves it.",
   22, HEAD, bold=True)
notes(s, "THE punchline slide. D is not a disappointment - it is the control that rules out "
         "the alternative explanation. Without D you could not claim form is what matters.")

# ---------------------------------------------------------------- 9 integrity
s = slide_head("Keeping the experiment clean", "worth saying out loud")
bullets(s, [
    "An earlier Model E scored 0.672  -  HIGHER than the one in the paper",
    "*We threw it out.",
    "It had also blended seasonality into the objective in three regions,",
    "so it changed the form AND the criterion at once. Two levers, not one.",
    "",
    "*Rebuilt holding the criterion fixed. Took 0.665 instead.",
], top=2.05, gap=0.55)
tb(s, 0.9, 5.85, 11.8, 0.8,
   "Gave up 0.008 of score to keep the attribution honest.", 21, GREEN, bold=True)
notes(s, "Lead with this if he questions rigour. It shows discipline: we gave up a better number "
         "to keep the experiment interpretable.")

# ---------------------------------------------------------------- 10 caveats
s = slide_head("Caveats to have ready", "do not get caught out")
bullets(s, [
    "*Magnitude:  E is 1.03x GFED5 globally, but that is largely COMPENSATING error",
    "   +323 Mha over-prediction cancelling -283 under. Net is excellent, regional detail less so.",
    "   Total absolute error is still the smallest of the four models.",
    "",
    "*Out of sample:  both held-out tests PASS",
    "   unseen years  r drops 0.05      unseen cells (blocked spatial CV)  r drops 0.015",
    "",
    "*Known weak spot:  seasonal timing. E peaks boreal Eurasia in October, GFED5 in April.",
], top=2.0, gap=0.52, size=18)
notes(s, "Do not oversell 1.03x. Say compensating error before he finds it. "
         "The held-out results are strong - lead with those if he presses on overfitting.")

# ---------------------------------------------------------------- 11 backup
s = slide_head("Backup  -  where things live", "if he asks")
bullets(s, [
    "Criterion (D):   spatial_taylor() in scripts/optimize_modelC_coupled.py,  SPATIAL_OBJ=1",
    "                 standalone scorer:  scripts/score_spatial.py",
    "Form (E):        fuel term + tropical suppression in fire_C, scripts/reproduce_modelC.py",
    "                 assembly:  scripts/assemble_continental.py  (clean preset)",
    "Params:          models/C/params.paperD.k1.json   |   models/E-clean/",
    "Output:          ilamb/MODELS_LEADERBOARD/ED-ModelC-E-clean/burntArea.nc",
], top=2.05, gap=0.55, size=17)
tb(s, 0.9, 5.7, 11.8, 0.9,
   "Trap:  the \"objective\" field inside the params JSON is a hardcoded label and still\n"
   "says \"ILAMB Overall\". It does NOT record the run objective. topk.paperD.json is authoritative.",
   17, ACCENT, bold=True)
notes(s, "Flag the hardcoded-label bug BEFORE he opens a file and sees the wrong objective string.")

out = "ModelCDE_advisor_deck.pptx"
prs.save(out)
print(f"wrote {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
